#!/usr/bin/env python3
"""
NEXUS-7 // PowerPoint Slide Generator
Builds a Blade Runner-themed .pptx for Tim Spann's 30-minute talk:
"Snowflake Vision: Beyond the Prompt — Enterprise AI"
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# ── Theme colours ─────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x0A, 0x14)
CYAN      = RGBColor(0x00, 0xFF, 0xF9)
MAGENTA   = RGBColor(0xFF, 0x00, 0xFF)
ORANGE    = RGBColor(0xFF, 0x66, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x88, 0x88, 0x88)
DARK_GREY = RGBColor(0x33, 0x33, 0x44)
CODE_BG   = RGBColor(0x12, 0x12, 0x20)


def _set_slide_bg(slide, color=BG_DARK):
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Consolas"):
    """Helper to add a styled textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_multiline(slide, left, top, width, height, lines, font_size=14,
                   color=WHITE, font_name="Consolas"):
    """Add a textbox with multiple lines, each as a separate paragraph."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
    return tf


def _add_code_block(slide, left, top, width, height, code, font_size=12):
    """Add a code block with dark background."""
    # Background rectangle
    from pptx.util import Inches as I
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = CYAN
    shape.line.width = Pt(1)
    # Text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = CYAN
        p.font.name = "Consolas"


def _add_table_slide(slide, left, top, width, rows_data, col_widths=None):
    """Add a styled table to a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top),
        Inches(width), Inches(0.4 * n_rows)
    )
    table = table_shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Consolas"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = CYAN
                else:
                    paragraph.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_DARK if r > 0 else DARK_GREY
    return table


# ══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.4, 9, 0.5,
                 "TYRELL CORPORATION // NEXUS DIVISION // 2026",
                 font_size=12, color=MAGENTA, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 1.2, 9, 1.2, "BEYOND THE PROMPT",
                 font_size=44, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 2.5, 9, 0.6, "Enterprise AI",
                 font_size=28, color=MAGENTA, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 3.6, 9, 0.8,
                 "Snowflake Vision: The Future of the AI Data Cloud",
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 4.8, 9, 0.5,
                 "Tim Spann  |  Sr Solution Engineer  |  Snowflake",
                 font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 5.5, 9, 0.5,
                 "NYC AI Meetup  |  March 2026",
                 font_size=14, color=GREY, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 6.5, 9, 0.4,
                 '"MORE HUMAN THAN HUMAN"',
                 font_size=11, color=DARK_GREY, alignment=PP_ALIGN.CENTER)


def slide_about(prs):
    """Slide 2 — About the Operator."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "ABOUT THE OPERATOR",
                 font_size=30, color=CYAN, bold=True)
    lines = [
        "DESIGNATION:     Tim Spann",
        "ROLE:            Sr Solution Engineer, Snowflake",
        "EXPERIENCE:      20+ years real-time data architectures",
        "CLEARANCE:       Apache Committer, Open Source Advocate",
        "SPECIALIZATION:  AI/ML, Streaming, IoT, Data Eng",
        "",
        "SIGNAL:          @PaaSDev",
        "UPLINK:          datainmotion.dev",
    ]
    _add_multiline(slide, 1.0, 1.3, 8, 4.5, lines, font_size=16, color=WHITE)


def slide_problem(prs):
    """Slide 3 — The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "THE REPLICANT PROBLEM",
                 font_size=30, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 1.0, 9, 0.5,
                 "Everyone has a chatbot. Few have enterprise AI.",
                 font_size=20, color=MAGENTA)
    data = [
        ["BASELINE CHATBOT", "NEXUS-7 ENTERPRISE AI"],
        ["Generic memory banks", "YOUR data matrix"],
        ["No governance", "Full RBAC mesh"],
        ["Isolated processing", "Integrated neural grid"],
        ["Prompt -> Response", "Context -> Insight -> Action"],
        ["No audit capability", "Complete lineage tracking"],
    ]
    _add_table_slide(slide, 0.8, 1.8, 8.4, data)
    _add_multiline(slide, 0.8, 4.8, 8.4, 2.0, [
        "The Gap: Demos work. Production fails. Why?",
        "  - AI doesn't understand YOUR business",
        "  - No governance = no trust = no deployment",
        "  - Isolated systems = fragmented intelligence",
    ], font_size=14, color=ORANGE)


def slide_vision(prs):
    """Slide 4 — AI Data Cloud Vision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "THE AI DATA CLOUD VISION",
                 font_size=30, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 1.0, 9, 0.5, "One Platform. Unified Intelligence.",
                 font_size=20, color=MAGENTA)
    lines = [
        "DATA MATRIX  ====>  CORTEX AI ENGINE  ====>  APPLICATIONS",
        "                                              (Streamlit)",
        "",
        "         UNIFIED GOVERNANCE & SECURITY PROTOCOLS",
        "            RBAC // AUDITING // ENCRYPTION",
    ]
    _add_code_block(slide, 0.8, 1.8, 8.4, 2.5, "\n".join(lines), font_size=14)
    _add_textbox(slide, 0.8, 4.8, 8.4, 0.5,
                 "AI processing happens WHERE your data lives",
                 font_size=16, color=ORANGE, bold=True)


def slide_four_cores(prs):
    """Slide 5 — Four Cores Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "FOUR REPLICANT CORES",
                 font_size=30, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.5, "The Enterprise AI Architecture",
                 font_size=18, color=MAGENTA)
    data = [
        ["CORE", "TECHNOLOGY", "FUNCTION"],
        ["NEURAL LAYER", "Cortex AI Functions", "AI primitives in SQL"],
        ["ENGRAM LAYER", "Semantic Views", "Business meaning for AI"],
        ["REASONING LAYER", "Cortex Agents", "Autonomous orchestration"],
        ["MESH LAYER", "MCP Protocol", "Universal integration"],
    ]
    _add_table_slide(slide, 0.8, 1.6, 8.4, data)
    lines = [
        "       APPLICATION INTERFACE",
        "              |",
        "       CORTEX AGENTS (REASONING)",
        "              |",
        "       SEMANTIC VIEWS (ENGRAMS)",
        "              |",
        "       CORTEX AI FUNCTIONS (NEURAL)",
    ]
    _add_code_block(slide, 1.5, 4.2, 7.0, 2.5, "\n".join(lines), font_size=13)


def slide_cortex_functions(prs):
    """Slide 6 — Core 1: Cortex AI Functions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORE 1: CORTEX AI FUNCTIONS",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4, "Neural Processing Units in SQL",
                 font_size=18, color=MAGENTA)
    code = """-- SENTIMENT ANALYSIS
SELECT AI_SENTIMENT('Flight experiencing turbulence')
    as emotional_state;

-- ZERO-SHOT CLASSIFICATION
SELECT AI_CLASSIFY(
    'Aircraft squawk 7700: Emergency',
    ['normal','radio_failure','hijack','emergency']
) as threat_classification;

-- AI COMPLETION (LLM)
SELECT AI_COMPLETE('claude-3-5-sonnet',
    'Summarize the operational status...'
) as synthesis;"""
    _add_code_block(slide, 0.5, 1.5, 9.0, 4.2, code, font_size=12)
    _add_textbox(slide, 0.5, 5.9, 9, 0.5,
                 "KEY: AI becomes a SQL operator. No data movement. Full governance.",
                 font_size=14, color=ORANGE, bold=True)


def slide_cortex_arsenal(prs):
    """Slide 7 — Full Cortex Function Arsenal."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORTEX AI FUNCTIONS — FULL ARSENAL",
                 font_size=28, color=CYAN, bold=True)
    data = [
        ["FUNCTION", "DESIGNATION", "USE CASE"],
        ["AI_COMPLETE", "Neural Synthesis", "Text generation, summaries"],
        ["AI_SENTIMENT", "Emotional Baseline", "Feedback analysis"],
        ["AI_CLASSIFY", "Pattern Recognition", "Zero-shot categorization"],
        ["AI_TRANSLATE", "Language Bridge", "Multi-language support"],
        ["AI_EMBED", "Vector Encoding", "Semantic search"],
        ["AI_PARSE_DOCUMENT", "Document Intel", "PDF/image extraction"],
    ]
    _add_table_slide(slide, 0.5, 1.1, 9.0, data)
    _add_multiline(slide, 0.8, 4.4, 8.4, 2.0, [
        "All functions execute with:",
        "  + Zero data migration",
        "  + Full RBAC enforcement",
        "  + Complete audit trail",
        "  + Warehouse-scale processing",
    ], font_size=14, color=CYAN)


def slide_semantic_views(prs):
    """Slide 8 — Core 2: Semantic Views."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORE 2: SEMANTIC VIEWS",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4, "The Missing Engram Layer",
                 font_size=18, color=MAGENTA)
    _add_multiline(slide, 0.5, 1.5, 9, 1.2, [
        'PROBLEM: "What\'s our revenue?"',
        'AI: ??? (Gross? Net? ARR? MRR? Which table?)',
        '',
        'SOLUTION: Semantic Views encode business meaning',
    ], font_size=15, color=WHITE)
    code = """-- SEMANTIC VIEW DEFINITION
dimensions:
  - name: customer_tier
    synonyms: [customer type, tier, segment]

metrics:
  - name: total_revenue
    expression: SUM(amount) WHERE status = 'COMPLETED'"""
    _add_code_block(slide, 0.5, 3.2, 9.0, 3.0, code, font_size=13)
    _add_textbox(slide, 0.5, 6.4, 9, 0.4,
                 "NO HALLUCINATION — Constrained to defined metrics",
                 font_size=14, color=ORANGE, bold=True)


def slide_semantic_arch(prs):
    """Slide 9 — Semantic View Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "SEMANTIC VIEW ARCHITECTURE",
                 font_size=28, color=CYAN, bold=True)
    lines = [
        'USER QUERY: "Show me revenue by customer tier"',
        "              |",
        "              v",
        "  SEMANTIC VIEW",
        "    Dimensions: customer_tier (Gold/Silver/Bronze)",
        "    Metrics: total_revenue = SUM(amount)",
        "    Synonyms: tier -> customer_tier",
        "              |",
        "              v",
        "  GENERATED SQL:",
        "    SELECT customer_tier, SUM(amount) as total_revenue",
        "    FROM orders WHERE status = 'COMPLETED'",
        "    GROUP BY customer_tier",
    ]
    _add_code_block(slide, 0.5, 1.0, 9.0, 5.5, "\n".join(lines), font_size=13)


def slide_cortex_agents(prs):
    """Slide 10 — Core 3: Cortex Agents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORE 3: CORTEX AGENTS",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4,
                 "Autonomous Reasoning Units",
                 font_size=18, color=MAGENTA)
    lines = [
        'USER: "Alert me if any flight drops below 1000ft near JFK"',
        "",
        "AGENT REASONING CHAIN:",
        "  1. INTENT ANALYSIS:   monitoring request detected",
        "  2. TOOL SELECTION:    flight_tracking_semantic_view",
        "  3. QUERY GENERATION:  flights near JFK, altitude < 1000",
        "  4. ACTION EXECUTION:  create notification rule",
        '  5. CONFIRMATION:      "Monitoring active..."',
    ]
    _add_code_block(slide, 0.5, 1.5, 9.0, 3.5, "\n".join(lines), font_size=13)
    _add_multiline(slide, 0.8, 5.3, 8.4, 1.8, [
        "CAPABILITIES:",
        "  + Tool selection & orchestration",
        "  + Context persistence across turns",
        "  + Semantic view integration",
        "  + Action execution (not just answers)",
    ], font_size=14, color=CYAN)


def slide_agent_arch(prs):
    """Slide 11 — Agent Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORTEX AGENT ARCHITECTURE",
                 font_size=28, color=CYAN, bold=True)
    lines = [
        "+--------------------------------------------------+",
        "|              CORTEX AGENT                         |",
        "|                                                   |",
        "| REASONING ----> TOOLS -------> ACTION EXECUTOR    |",
        "|  ENGINE         - Semantic      |                 |",
        "|    ^               Views        v                 |",
        "|    |             - Cortex     RESPONSE             |",
        "| CONTEXT           Search                          |",
        "| MEMORY          - SQL Exec                        |",
        "|                 - Custom Fns                      |",
        "+--------------------------------------------------+",
    ]
    _add_code_block(slide, 0.5, 1.0, 9.0, 4.5, "\n".join(lines), font_size=13)
    _add_textbox(slide, 0.5, 5.8, 9, 0.5,
                 "35 agents discovered in demo account",
                 font_size=16, color=ORANGE, bold=True)


def slide_mcp(prs):
    """Slide 12 — Core 4: MCP Protocol."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "CORE 4: MCP PROTOCOL",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4,
                 "Model Context Protocol — Universal Mesh Integration",
                 font_size=18, color=MAGENTA)
    _add_multiline(slide, 0.5, 1.6, 9, 1.5, [
        "WHAT IS MCP?",
        "  - Open protocol for AI-to-tool communication",
        "  - Standardized interface for external system access",
        "  - Snowflake provides managed MCP server",
    ], font_size=15, color=WHITE)
    code = """# CONNECT SNOWFLAKE TO EXTERNAL AI
cortex mcp add snowflake https://account.snowflakecomputing.com/mcp

# NOW ANY LLM CAN ACCESS:
#   - Your Snowflake data
#   - Your semantic views
#   - Your Cortex agents
#   - With full governance"""
    _add_code_block(slide, 0.5, 3.2, 9.0, 2.8, code, font_size=13)
    _add_multiline(slide, 0.5, 6.2, 9, 0.8, [
        "WHY: Claude, GPT, any LLM -> Snowflake data with auth & audit",
    ], font_size=14, color=ORANGE)


def slide_demo_intro(prs):
    """Slide 13 — Live Demo intro."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.5, 9, 0.8, "LIVE DEMO",
                 font_size=44, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 1.5, 9, 0.5, "NEXUS-7 IN ACTION",
                 font_size=24, color=MAGENTA, alignment=PP_ALIGN.CENTER)
    lines = [
        "CORTEX AI FUNCTIONS    [ONLINE]   Neural SQL operators",
        "SEMANTIC VIEWS         [ONLINE]   45 views discovered",
        "CORTEX AGENTS          [ONLINE]   35 agents available",
        "MCP INTEGRATION        [ONLINE]   Mesh connected",
        "",
        "DATA DOMAINS:",
        "  + Real-time flight tracking (ADS-B)",
        "  + Financial market data",
        "  + IoT sensor networks",
    ]
    _add_code_block(slide, 0.8, 2.5, 8.4, 4.0, "\n".join(lines), font_size=14)


def slide_demo_cortex(prs):
    """Slide 14 — Demo: Cortex AI Functions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "DEMO 1: CORTEX AI FUNCTIONS",
                 font_size=28, color=CYAN, bold=True)
    code = """-- SENTIMENT
SELECT AI_SENTIMENT('Flight AA123 turbulence, passengers secure');
-- Result: 0.15 (neutral)

-- CLASSIFICATION
SELECT AI_CLASSIFY(
    'Aircraft squawk 7700: Emergency declared',
    ['normal', 'radio_failure', 'hijack', 'emergency']
);
-- Result: {"label": "emergency", "score": 0.95}

-- SYNTHESIS
SELECT AI_COMPLETE('claude-3-5-sonnet',
    'Summarize: Boeing 737 hydraulic failure, landed safely');
-- Result: "A Boeing 737 experienced hydraulic warnings
--          and safely diverted after declaring emergency." """
    _add_code_block(slide, 0.3, 1.0, 9.4, 5.5, code, font_size=12)


def slide_demo_semantic(prs):
    """Slide 15 — Demo: Semantic View Queries."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "DEMO 2: SEMANTIC VIEW QUERIES",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4,
                 "Natural Language -> Precise SQL",
                 font_size=18, color=MAGENTA)
    data = [
        ["METRIC", "VALUE", "DESCRIPTION"],
        ["TOTAL_AIRCRAFT", "26", "Distinct aircraft in grid"],
        ["TOTAL_FLIGHTS", "21", "Active flight vectors"],
        ["AVG_ALTITUDE", "19,137 ft", "Mean cruising altitude"],
        ["AVG_GROUND_SPEED", "330 kts", "Mean velocity"],
    ]
    _add_table_slide(slide, 0.5, 1.5, 9.0, data)
    code = """-- GENERATED SQL (via semantic view)
SELECT sv.total_aircraft, sv.avg_altitude, sv.avg_ground_speed
FROM SEMANTIC_VIEW(DEMO.DEMO.adsb_flight_tracking
    METRICS total_aircraft, avg_altitude, avg_ground_speed) AS sv;"""
    _add_code_block(slide, 0.5, 4.0, 9.0, 2.0, code, font_size=12)


def slide_demo_agent(prs):
    """Slide 16 — Demo: Agent Orchestration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "DEMO 3: CORTEX AGENT ORCHESTRATION",
                 font_size=26, color=CYAN, bold=True)
    lines = [
        'USER: "Compare flight activity with weather conditions"',
        "",
        "AGENT EXECUTION:",
        "  STEP 1: Intent Analysis",
        "    -> Cross-domain correlation request",
        "  STEP 2: Tool Selection",
        "    -> ADSB_FLIGHT_TRACKING semantic view",
        "    -> AIRQUALITY semantic view",
        "  STEP 3: Parallel Query Execution",
        "    -> Flight: 26 aircraft, 330 kts avg",
        "    -> Weather: AQI, visibility",
        "  STEP 4: Correlation & Synthesis",
        '  STEP 5: "Current airspace shows normal activity..."',
    ]
    _add_code_block(slide, 0.3, 1.0, 9.4, 5.5, "\n".join(lines), font_size=13)


def slide_demo_oracle(prs):
    """Slide 17 — Demo: AI-Generated Insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "DEMO 4: THE ORACLE FUNCTION",
                 font_size=28, color=CYAN, bold=True)
    code = """SELECT AI_COMPLETE(
    'claude-3-5-sonnet',
    'You are NEXUS-7, an advanced aviation AI.
     Based on this telemetry, provide operational insight:
     - Total aircraft: 26
     - Active flights: 21
     - Mean altitude: 19,137 feet
     - Mean velocity: 330 knots'
) as oracle_insight;"""
    _add_code_block(slide, 0.5, 1.0, 9.0, 3.2, code, font_size=13)
    _add_textbox(slide, 0.5, 4.5, 9, 0.3, "ORACLE OUTPUT:", font_size=14,
                 color=MAGENTA, bold=True)
    _add_multiline(slide, 0.5, 4.9, 9, 2.2, [
        '"Based on the telemetry data, the current airspace shows',
        'moderate traffic density with 26 total aircraft, of which',
        '80.7% are maintaining active flight vectors, suggesting',
        'normal en-route operations. The mean altitude of 19,137',
        'feet combined with average velocity of 330 knots indicates',
        'most aircraft are in the mid-flight cruise phase."',
    ], font_size=13, color=WHITE)


def slide_enterprise_diff(prs):
    """Slide 18 — Enterprise AI Difference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "THE ENTERPRISE AI DIFFERENCE",
                 font_size=28, color=CYAN, bold=True)
    data = [
        ["BASELINE CHATBOT", "NEXUS-7 ENTERPRISE AI"],
        ["Generic training data", "YOUR data matrix"],
        ["No access control", "Full RBAC mesh"],
        ["Isolated responses", "Integrated workflows"],
        ["Prompt -> Response", "Context -> Insight -> Action"],
        ["No accountability", "Complete audit trail"],
        ["Demo-only", "Production-ready"],
    ]
    _add_table_slide(slide, 0.8, 1.0, 8.4, data)
    _add_textbox(slide, 0.5, 5.2, 9, 0.8,
                 '"A SYSTEM WITHOUT GOVERNANCE IS JUST AN EXPENSIVE TOY"',
                 font_size=18, color=MAGENTA, bold=True,
                 alignment=PP_ALIGN.CENTER)


def slide_architecture_pattern(prs):
    """Slide 19 — AI Data Product Pattern."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "ARCHITECTURE: AI DATA PRODUCT",
                 font_size=28, color=CYAN, bold=True)
    lines = [
        "APPLICATION LAYER",
        "  Streamlit App | REST API | Cortex Code CLI",
        "              |",
        "CORTEX AGENT",
        "  Tools: [Semantic Views, Search, Functions]",
        "              |",
        "SEMANTIC LAYER",
        "  Semantic Views defining business metrics",
        "              |",
        "DATA LAYER",
        "  Tables | Dynamic Tables | Streams | Iceberg",
    ]
    _add_code_block(slide, 1.0, 1.0, 8.0, 5.0, "\n".join(lines), font_size=14)


def slide_deployments(prs):
    """Slide 20 — Real-World Deployments."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "REAL-WORLD DEPLOYMENTS",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 1.0, 4, 0.4, "Aviation & Transportation",
                 font_size=18, color=MAGENTA, bold=True)
    _add_multiline(slide, 0.5, 1.5, 4, 1.5, [
        "  Real-time flight monitoring",
        "  Emergency squawk detection",
        "  Weather impact correlation",
    ], font_size=14, color=WHITE)
    _add_textbox(slide, 0.5, 3.0, 4, 0.4, "Financial Services",
                 font_size=18, color=MAGENTA, bold=True)
    _add_multiline(slide, 0.5, 3.5, 4, 1.5, [
        "  Market sentiment analysis",
        "  Document intelligence",
        "  Regulatory compliance",
    ], font_size=14, color=WHITE)
    _add_textbox(slide, 5.0, 1.0, 5, 0.4, "IoT & Manufacturing",
                 font_size=18, color=MAGENTA, bold=True)
    _add_multiline(slide, 5.0, 1.5, 5, 1.5, [
        "  Sensor anomaly detection",
        "  Predictive maintenance",
        "  Supply chain optimization",
    ], font_size=14, color=WHITE)


def slide_getting_started(prs):
    """Slide 21 — Getting Started."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "GETTING STARTED",
                 font_size=28, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 0.9, 9, 0.4,
                 "Your Path to Enterprise AI",
                 font_size=18, color=MAGENTA)
    code = """-- STEP 1: Define your semantic layer
CREATE SEMANTIC VIEW my_business_metrics AS ...

-- STEP 2: Build your Cortex Agent
CREATE CORTEX AGENT my_analyst
  WITH TOOLS = (my_semantic_view, cortex_search_service)

-- STEP 3: Deploy with Cortex Code
cortex agent query "What were sales last quarter?" \\
  --agent=my_analyst"""
    _add_code_block(slide, 0.5, 1.5, 9.0, 4.0, code, font_size=14)


def slide_takeaways(prs):
    """Slide 22 — Key Takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "KEY TAKEAWAYS",
                 font_size=30, color=CYAN, bold=True)
    _add_textbox(slide, 0.5, 1.0, 9, 0.4,
                 "Beyond the Prompt: Enterprise AI", font_size=20, color=MAGENTA)
    items = [
        ("1.", "SEMANTIC VIEWS", "AI that understands YOUR business"),
        ("2.", "CORTEX AGENTS", "Reasoning + orchestration + action"),
        ("3.", "MCP PROTOCOL", "Universal tool integration"),
        ("4.", "CORTEX CODE", "Developer-first AI experience"),
        ("5.", "GOVERNANCE", "Security you can trust"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 1.7 + i * 0.8
        _add_textbox(slide, 0.8, y, 0.5, 0.5, num,
                     font_size=18, color=ORANGE, bold=True)
        _add_textbox(slide, 1.4, y, 3, 0.5, title,
                     font_size=18, color=CYAN, bold=True)
        _add_textbox(slide, 4.5, y, 5.5, 0.5, desc,
                     font_size=16, color=WHITE)
    _add_textbox(slide, 0.5, 6.0, 9, 0.5,
                 "Enterprise AI that works in production, not just in demos.",
                 font_size=16, color=ORANGE, bold=True,
                 alignment=PP_ALIGN.CENTER)


def slide_resources(prs):
    """Slide 23 — Resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "RESOURCES",
                 font_size=30, color=CYAN, bold=True)
    lines = [
        "DOCUMENTATION",
        "  docs.snowflake.com/cortex",
        "",
        "QUICKSTARTS",
        "  quickstarts.snowflake.com",
        "",
        "COMMUNITY",
        "  community.snowflake.com",
        "",
        "SOURCE CODE",
        "  github.com/Snowflake-Labs",
    ]
    _add_code_block(slide, 1.0, 1.0, 8.0, 4.5, "\n".join(lines), font_size=15)
    _add_multiline(slide, 1.0, 5.8, 8, 1.0, [
        "OPERATOR CONTACT:  @PaaSDev  |  datainmotion.dev",
    ], font_size=16, color=ORANGE)


def slide_qa(prs):
    """Slide 24 — Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 1.5, 9, 1.0, "Q & A",
                 font_size=54, color=CYAN, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 3.0, 9, 0.5, "Questions?",
                 font_size=24, color=MAGENTA, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 4.5, 9, 0.5,
                 "Tim Spann  |  @PaaSDev  |  datainmotion.dev",
                 font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, 0.5, 6.2, 9, 0.4,
                 '"MORE HUMAN THAN HUMAN"',
                 font_size=12, color=DARK_GREY, alignment=PP_ALIGN.CENTER)


def slide_appendix_commands(prs):
    """Slide 25 — Appendix: Command Reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "APPENDIX: COMMAND REFERENCE",
                 font_size=24, color=CYAN, bold=True)
    code = """# DISCOVER SEMANTIC VIEWS
cortex semantic-views discover
cortex semantic-views describe DEMO.DEMO.ADSB_FLIGHT_TRACKING

# QUERY VIA CORTEX ANALYST
cortex analyst query "average altitude of flights" \\
  --view=DEMO.DEMO.ADSB_FLIGHT_TRACKING

# DISCOVER AGENTS
cortex agents discover
cortex agents describe DEMO.DEMO.AVIATION_WEATHER_AGENT

# START STREAMLIT DEMO
./manage.sh start"""
    _add_code_block(slide, 0.3, 1.0, 9.4, 5.5, code, font_size=13)


def slide_validation(prs):
    """Slide 26 — Validation Status."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_textbox(slide, 0.5, 0.3, 9, 0.5, "VALIDATION STATUS",
                 font_size=28, color=CYAN, bold=True)
    lines = [
        "NEXUS-7 SYSTEM VALIDATION",
        "",
        "+ AI_SENTIMENT         OPERATIONAL",
        "+ AI_CLASSIFY          OPERATIONAL",
        "+ AI_COMPLETE          OPERATIONAL",
        "+ SEMANTIC_VIEW        OPERATIONAL (45 views)",
        "+ CORTEX_AGENTS        OPERATIONAL (35 agents)",
        "+ ADSB_TELEMETRY       26 aircraft, 21 flights",
        "+ FINANCIAL_DATA       Historical available",
        "+ STREAMLIT_APP        http://localhost:8509",
        "",
        "ALL SYSTEMS NOMINAL // READY FOR DEPLOYMENT",
    ]
    _add_code_block(slide, 0.5, 1.0, 9.0, 5.5, "\n".join(lines), font_size=14)


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def build_presentation(output_path=None):
    """Build the full slide deck and save to output_path."""
    if output_path is None:
        output_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)), "presentation.pptx"
        )

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build all slides in order
    slide_title(prs)              # 1
    slide_about(prs)              # 2
    slide_problem(prs)            # 3
    slide_vision(prs)             # 4
    slide_four_cores(prs)         # 5
    slide_cortex_functions(prs)   # 6
    slide_cortex_arsenal(prs)     # 7
    slide_semantic_views(prs)     # 8
    slide_semantic_arch(prs)      # 9
    slide_cortex_agents(prs)      # 10
    slide_agent_arch(prs)         # 11
    slide_mcp(prs)                # 12
    slide_demo_intro(prs)         # 13
    slide_demo_cortex(prs)        # 14
    slide_demo_semantic(prs)      # 15
    slide_demo_agent(prs)         # 16
    slide_demo_oracle(prs)        # 17
    slide_enterprise_diff(prs)    # 18
    slide_architecture_pattern(prs)  # 19
    slide_deployments(prs)        # 20
    slide_getting_started(prs)    # 21
    slide_takeaways(prs)          # 22
    slide_resources(prs)          # 23
    slide_qa(prs)                 # 24
    slide_appendix_commands(prs)  # 25
    slide_validation(prs)         # 26

    prs.save(output_path)
    print(f"Saved {len(prs.slides)} slides to {output_path}")
    return output_path


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else None
    build_presentation(out)
